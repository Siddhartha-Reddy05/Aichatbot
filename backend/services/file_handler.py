import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from typing import List, Tuple
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image
import logging
import io

logger = logging.getLogger(__name__)

def extract_text_from_file(filename: str, content: bytes) -> List[Tuple[str, int, str]]:
    ext = os.path.splitext(filename)[1].lower()
    
    print(f"📄 Processing file: {filename} with extension: {ext}")
    
    if ext == ".txt":
        return extract_text_from_txt(content, filename)
    elif ext == ".pdf":
        return extract_text_from_pdf_enhanced(content, filename)
    elif ext == ".docx":
        return extract_text_from_docx(content, filename)
    elif ext == ".pptx":
        return extract_text_from_pptx(content, filename)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def extract_text_from_pdf_enhanced(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    """Enhanced PDF text extraction with multiple fallback methods."""
    pages = []
    
    try:
        with fitz.open(stream=content, filetype="pdf") as doc:
            print(f"📊 PDF has {len(doc)} pages")
            
            for i, page in enumerate(doc):
                print(f"🔍 Processing page {i+1}...")
                
                # Method 1: Try standard text extraction
                text = page.get_text().strip()
                
                if text and len(text) > 50:  # Reasonable amount of text
                    print(f"✅ Page {i+1}: Found {len(text)} characters with standard extraction")
                    pages.append((text, i + 1, source))
                    continue
                
                # Method 2: Try "textpage" extraction (different method)
                try:
                    textpage = page.get_textpage()
                    text = textpage.extractText().strip()
                    if text and len(text) > 50:
                        print(f"✅ Page {i+1}: Found {len(text)} characters with textpage extraction")
                        pages.append((text, i + 1, source))
                        continue
                except:
                    pass
                
                # Method 3: Check if it's an image-based PDF and try OCR
                image_list = page.get_images()
                if image_list:
                    print(f"🖼️ Page {i+1}: Contains {len(image_list)} images, attempting OCR...")
                    try:
                        # Get page as image
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution for better OCR
                        img_data = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_data))
                        
                        # Perform OCR
                        ocr_text = pytesseract.image_to_string(img).strip()
                        
                        if ocr_text and len(ocr_text) > 20:
                            print(f"✅ Page {i+1}: Found {len(ocr_text)} characters with OCR")
                            pages.append((ocr_text, i + 1, source))
                        else:
                            print(f"⚠️ Page {i+1}: OCR found minimal or no text")
                            pages.append(("[Image-based page - minimal text extracted]", i + 1, source))
                    except Exception as ocr_error:
                        print(f"❌ Page {i+1}: OCR failed: {ocr_error}")
                        pages.append(("[Image-based page - OCR failed]", i + 1, source))
                else:
                    print(f"⚠️ Page {i+1}: No extractable text found")
                    pages.append(("[No text content]", i + 1, source))
            
    except Exception as e:
        print(f"❌ Error processing PDF: {e}")
        # Try fallback method: pdf2image + OCR
        try:
            print("🔄 Trying fallback PDF processing with pdf2image...")
            images = convert_from_bytes(content)
            for i, image in enumerate(images):
                text = pytesseract.image_to_string(image).strip()
                pages.append((text, i + 1, source))
                print(f"✅ Fallback: Page {i+1} extracted {len(text)} characters")
        except Exception as fallback_error:
            print(f"❌ Fallback also failed: {fallback_error}")
            raise ValueError(f"Failed to extract text from PDF: {e}")
    
    print(f"✅ PDF extraction complete. Found {len(pages)} pages.")
    return pages


def extract_text_from_docx(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    from io import BytesIO
    try:
        doc = docx.Document(BytesIO(content))
        full_text = "\n".join([para.text for para in doc.paragraphs])
        print(f"📝 DOCX text length: {len(full_text)} characters")
        return [(full_text, 1, source)]
    except Exception as e:
        print(f"❌ Error extracting text from DOCX: {e}")
        raise

def extract_text_from_pptx(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    from io import BytesIO
    try:
        prs = Presentation(BytesIO(content))
        pages = []
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            pages.append((slide_text.strip(), i + 1, source))
        print(f"📝 PPTX text extracted from {len(pages)} slides")
        return pages
    except Exception as e:
        print(f"❌ Error extracting text from PPTX: {e}")
        raise

def extract_text_from_txt(content: bytes, source: str) -> List[Tuple[str, int, str]]:
    try:
        text = content.decode("utf-8", errors="ignore")
        print(f"📝 TXT text length: {len(text)} characters")
        return [(text, 1, source)]
    except Exception as e:
        print(f"❌ Error extracting text from TXT: {e}")
        raise